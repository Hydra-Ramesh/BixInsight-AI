import sys
import json
import base64
import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_ppt(analysis_data):
    prs = Presentation()
    
    # Optional: Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    filename = analysis_data.get("fileName", "Data Analysis")
    title.text = f"BixInsight AI Report: {filename}"
    subtitle.text = f"Generated automatically by BixInsight AI\nRows: {analysis_data.get('rowCount', 0)} | Columns: {analysis_data.get('columnCount', 0)}"

    # Summary Slide
    summary = analysis_data.get("summary")
    if summary:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Executive Summary"
        tf = body_shape.text_frame
        tf.text = summary

    # Charts Slides
    charts = analysis_data.get("charts", [])
    for chart in charts:
        img_b64 = chart.get("imageBase64")
        if not img_b64:
            continue
            
        slide_layout = prs.slide_layouts[5] # Blank with title
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = chart.get("title", "Chart Analysis")
        
        # Recommendation / Description Text Box
        rec_text = chart.get("recommendation", chart.get("description", ""))
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.8))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = rec_text
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        
        # Decode image
        try:
            img_data = base64.b64decode(img_b64)
            img_stream = io.BytesIO(img_data)
            
            # Add image to slide
            left = Inches(1.5)
            top = Inches(2.0)
            height = Inches(5.0)
            pic = slide.shapes.add_picture(img_stream, left, top, height=height)
        except Exception as e:
            print(f"Failed to add chart image: {e}", file=sys.stderr)

    # Insights Slide
    insights = analysis_data.get("insights", [])
    if insights:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Key AI Insights"
        
        tf = slide.placeholders[1].text_frame
        tf.clear() # clear default paragraphs
        
        for ins in insights[:5]: # Top 5 to fit
            p = tf.add_paragraph()
            p.text = f"[{ins.get('importance', 'info').upper()}] {ins.get('title', 'Insight')}: {ins.get('description', '')}"
            p.font.size = Pt(14)
            
    # Save to memory
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()

if __name__ == "__main__":
    try:
        input_data = sys.stdin.read()
        analysis = json.loads(input_data)
        
        ppt_bytes = create_ppt(analysis)
        
        # Write bytes back out to stdout
        sys.stdout.buffer.write(ppt_bytes)
        
    except Exception as e:
        print(f"Error generating PPTX: {e}", file=sys.stderr)
        sys.exit(1)
